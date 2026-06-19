from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

prs = Presentation()

# Function to add a slide with title and bullet points
def add_slide(title_text, content_text):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    for i, line in enumerate(content_text):
        if i == 0:
            tf.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
        
# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "BandOps"
subtitle.text = "Multi-Agent Command Rooms for Regulated Crisis Operations.\nPowered by the Band API"

# Slide 2: The Problem
add_slide("The Problem: Aircraft on Ground (AOG)", [
    "A grounded plane costs airlines $15,000 per minute.",
    "Mechanics, Logistics, Flight Ops, and the FAA are siloed.",
    "Communication is chaotic: emails, phone calls, whiteboards.",
    "Result: Hours of delays, compliance risks, and millions lost."
])

# Slide 3: The Solution
add_slide("The Solution: BandOps Command Center", [
    "Replaces siloed human panic with an orchestrated swarm of specialized AI agents.",
    "Agents analyze faults, locate parts, and check regulations in seconds.",
    "Human-in-the-loop: Agents recommend, humans approve.",
    "Built for high-stakes, highly regulated enterprise environments."
])

# Slide 4: Why Band?
add_slide("Why Band? (The Interoperability Layer)", [
    "Band provides the persistent room, identity verification, and audit trail.",
    "Framework-agnostic: Allows LangGraph, CrewAI, and Python agents to collaborate securely.",
    "Provides the zero-trust mesh that makes enterprise-grade multi-agent systems possible."
])

# Slide 5: The War Room Demo
add_slide("The War Room Architecture", [
    "1. Line Mechanic diagnoses the fault.",
    "2. Supply Chain sources replacement parts.",
    "3. Flight Ops calculates delay costs.",
    "4. FAA Compliance acts as an adversarial auditor.",
    "5. AOG Commander synthesizes the debate for the human."
])

# Slide 6: Immutable Auditability
add_slide("Immutable Auditability", [
    "In aviation, every single decision must be auditable.",
    "BandOps logs the entire agent debate to an immutable ledger.",
    "Logs the final human approval cryptographically.",
    "Zero black boxes. Full traceability for regulators."
])

# Slide 7: Beyond Aviation
add_slide("Platform Scalability", [
    "BandOps architecture applies to any high-stakes, regulated environment:",
    "- Hospitals (ER Triage & Bed Management)",
    "- Energy Grids (Disaster Response & Load Balancing)",
    "- Global Logistics (Supply Chain Shock Recovery)"
])

# Slide 8: Thank You
add_slide("Thank You", [
    "Try the Demo locally or view the GitHub Repo.",
    "Built with love for the Band of Agents Hackathon."
])

os.makedirs("pitch", exist_ok=True)
prs.save("pitch/BandOps_Pitch_Deck.pptx")
print("PPTX generated successfully!")
